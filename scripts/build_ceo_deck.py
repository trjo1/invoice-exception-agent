"""Generate the CEO demo slide deck (.pptx) from python-pptx.

Uses the locked TruVs brand tokens (navy / red / green / soft band) and
Arial body. 12-slide structure matching the talking-points doc at
docs/presentations/ceo_demo_talking_points.md.

Run: `uv run python scripts/build_ceo_deck.py`
Output: `docs/presentations/ceo_demo_agent1.pptx`
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# --- Brand tokens (locked, from CLAUDE.md) ---
NAVY = RGBColor(0x1F, 0x3A, 0x68)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x1A, 0x6D, 0x3F)
SOFT = RGBColor(0xE8, 0xEE, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5C, 0x63, 0x70)
AMBER = RGBColor(0xF3, 0xC3, 0x4A)

FONT = "Arial"

OUTPUT = Path("docs/presentations/ceo_demo_agent1.pptx")


def _set_text(tf, text: str, *, size: int = 18, bold: bool = False,
              color: RGBColor = TEXT, align=PP_ALIGN.LEFT) -> None:
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.runs[0]
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_brand_strip(slide, title: str) -> None:
    """Top navy strip with white title — the signature TruVs slide header."""
    strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.85),
    )
    strip.fill.solid()
    strip.fill.fore_color.rgb = NAVY
    strip.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12.3), Inches(0.55))
    _set_text(tx.text_frame, title, size=26, bold=True, color=WHITE)


def _add_bullets(slide, bullets: list[str], *,
                 top: float = 1.4, left: float = 0.7,
                 width: float = 12.0, height: float = 5.5,
                 size: int = 18) -> None:
    """Bullet list."""
    tx = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    tf = tx.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        # python-pptx doesn't expose bullet style; we'll prefix with •
        for run in p.runs:
            run.font.name = FONT
            run.font.size = Pt(size)
            run.font.color.rgb = TEXT


def _add_caption(slide, text: str, *, top: float, color: RGBColor = MUTED) -> None:
    tx = slide.shapes.add_textbox(
        Inches(0.7), Inches(top), Inches(12.0), Inches(0.5),
    )
    _set_text(tx.text_frame, text, size=14, color=color)


def _add_footer(slide, page_no: int, total: int = 12) -> None:
    tx = slide.shapes.add_textbox(
        Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.4),
    )
    _set_text(tx.text_frame, f"{page_no} / {total}",
              size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def _add_callout_box(slide, text: str, *, top: float, left: float = 0.7,
                     width: float = 12.0, height: float = 1.2,
                     bg: RGBColor = SOFT) -> None:
    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = bg
    box.line.color.rgb = NAVY
    box.line.width = Pt(0.75)
    box.text_frame.margin_left = Inches(0.2)
    box.text_frame.margin_top = Inches(0.15)
    box.text_frame.margin_right = Inches(0.2)
    _set_text(box.text_frame, text, size=15, color=NAVY, bold=False)


def build() -> Presentation:
    pres = Presentation()
    pres.slide_width = Inches(13.33)   # 16:9
    pres.slide_height = Inches(7.5)
    blank = pres.slide_layouts[6]

    # ============== Slide 1 — Cover ==============
    s = pres.slides.add_slide(blank)
    # Full navy background
    bg = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5),
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    # Title
    tx = s.shapes.add_textbox(Inches(0.7), Inches(2.6), Inches(12.0), Inches(1.2))
    _set_text(tx.text_frame, "Agent 1: P2P Exception Orchestrator",
              size=44, bold=True, color=WHITE)
    # Subtitle
    tx = s.shapes.add_textbox(Inches(0.7), Inches(3.9), Inches(12.0), Inches(0.7))
    _set_text(tx.text_frame, "Reference implementation of the Process-to-Agent Method",
              size=22, color=SOFT)
    # Date strip
    tx = s.shapes.add_textbox(Inches(0.7), Inches(6.5), Inches(12.0), Inches(0.4))
    _set_text(tx.text_frame, "TruVs AI Practice · TJ · 2026-05-13",
              size=14, color=SOFT)

    # ============== Slide 2 — The pain ==============
    s = pres.slides.add_slide(blank)
    _add_brand_strip(s, "The pain we're attacking")
    _add_caption(s,
                 "Fortune-500 P2P operations drown in exception handling — it's the single biggest "
                 "drag on AP throughput.", top=1.05)
    _add_bullets(s, [
        "•  30–40% of AP team time spent on exception handling",
        "•  5–15 day working-capital cycle drag from coordination delays",
        "•  Real-money leakage: paid duplicates, missed early-pay discounts, supplier friction",
        "•  At 50K invoices/month, ~15–20% of every invoice has something wrong",
        "•  AP supervisor → procurement → buyer → supplier email chains. Days per exception.",
    ], top=1.7, size=18)
    _add_callout_box(s,
                     "Two-thirds of an AP team's time burns on coordination — not on processing.",
                     top=5.7, height=0.9)
    _add_footer(s, 2)

    # ============== Slide 3 — Why existing tools don't solve it ==============
    s = pres.slides.add_slide(blank)
    _add_brand_strip(s, "Why existing tools don't solve this")
    _add_caption(s,
                 "Every buyer has tried one or more of these. None hold the end-to-end exception lifecycle.",
                 top=1.05)
    _add_bullets(s, [
        "•  Document AI (Botminds, AI Builder, ABBYY) — extracts fields. Doesn't decide.",
        "•  ERP workflow (SAP S/4, Ariba Approvals) — handles structured approvals. Can't reason on unstructured supplier comms.",
        "•  RPA (UiPath, Blue Prism) — follows rules. Breaks on long-tail exceptions, which are 100% of the problem.",
        "•  Generic chatbots (Decagon, ChatGPT) — don't hold workflow state. Can't write back to SAP.",
    ], top=1.7, size=18)
    _add_callout_box(s,
                     "The problem isn't extraction / workflow / rules. The problem is COORDINATION "
                     "across systems with ambiguous facts and varied resolutions. That's our gap.",
                     top=5.5, height=1.3)
    _add_footer(s, 3)

    # ============== Slide 4 — What we built ==============
    s = pres.slides.add_slide(blank)
    _add_brand_strip(s, "What we built — 9-node exception orchestrator")
    _add_caption(s,
                 "Reference implementation of Pattern 2 (Multi-System Process Agent). Shipped end-to-end.",
                 top=1.05)
    _add_bullets(s, [
        "1.  Extract  →  PDF → structured JSON (invoice fields, line items, tax, total)",
        "2.  Cross-case context  →  Vendor master + PO + GR + invoice history lookups",
        "3.  Classify  →  One of 13 exception categories with confidence + evidence",
        "4.  Retrieve  →  Top-5 relevant policy snippets via local embeddings",
        "5.  Decide  →  Recommended action + rationale + counterfactual",
        "6.  Route  →  Pure-rules 3-tier HITL routing (16-action table)",
        "7.  Draft  →  Supplier email or internal note (never auto-sent)",
        "8.  Approval queue  →  SQLite-backed inbox + FastAPI demo console",
        "9.  Execute  →  Mock today; real SAP/Ariba/ServiceNow write-back pending creds",
    ], top=1.6, size=15)
    _add_callout_box(s,
                     "Headline: 9/9 nodes shipped · 24 golden cases · 80 tests · Total LLM spend on the build: $3.74",
                     top=6.4, height=0.9, bg=SOFT)
    _add_footer(s, 4)

    # ============== Slide 5 — Why this is the TruVs wedge ==============
    s = pres.slides.add_slide(blank)
    _add_brand_strip(s, "Why this is the TruVs wedge")
    _add_caption(s,
                 "Sits at the intersection of Process-First × Build — the Accenture quadrant. "
                 "Reference implementation that anchors every P2A engagement.",
                 top=1.05)
    _add_bullets(s, [
        "•  Decision Framework: Process-First × Build  ⟶  this agent's home quadrant",
        "•  Maps to the P2A flagship offering ($600K – $2.5M engagement range)",
        "•  $200K build floor (locked) — no sub-$200K engagements",
        "•  Stage 3 kill gate: we walk away from use cases that don't pencil out",
        "•  The wedge: \"the firm that says NO to bad AI use cases — proves it with the process data\"",
    ], top=1.8, size=18)
    _add_callout_box(s,
                     "Every paid engagement starts FROM this reference implementation — not from scratch. "
                     "That's the cost-of-build leverage.",
                     top=5.7, height=1.0)
    _add_footer(s, 5)

    # ============== Slide 6 — Demo intro ==============
    s = pres.slides.add_slide(blank)
    _add_brand_strip(s, "Live demo — three things to watch for")
    _add_caption(s, "Switching to a browser. Watch the agent reason about a real invoice.",
                 top=1.05)
    # Three big numbered items
    items = [
        ("1", "Visibility — every step is visible. No black box.",
         "You see what it extracted, classified as, retrieved, and recommended — with rationale."),
        ("2", "Counterfactual — \"if X were different, the answer would be Y.\"",
         "Builds trust faster than accuracy claims alone. The reviewer can verify reasoning."),
        ("3", "Human-in-the-loop — money-moving actions never auto-send.",
         "Tier 2 review for supplier comms / PO amendments. Tier 3 for fraud / treasury / VP finance."),
    ]
    y = 1.7
    for num, headline, sub in items:
        # Number badge
        badge = s.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.7), Inches(y), Inches(0.7), Inches(0.7),
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = NAVY
        badge.line.fill.background()
        _set_text(badge.text_frame, num, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Text
        tx = s.shapes.add_textbox(Inches(1.6), Inches(y), Inches(11.0), Inches(0.5))
        _set_text(tx.text_frame, headline, size=18, bold=True, color=NAVY)
        tx2 = s.shapes.add_textbox(Inches(1.6), Inches(y + 0.5), Inches(11.0), Inches(0.6))
        _set_text(tx2.text_frame, sub, size=14, color=MUTED)
        y += 1.4
    _add_footer(s, 6)

    # ============== Slide 7 — How it works under the hood ==============
    # (Slide numbering matches the talking-points doc, which has slide 7 as
    # "post-demo recap, the 9 nodes." We're past the demo here.)
    s = pres.slides.add_slide(blank)
    _add_brand_strip(s, "How it works — under the hood")
    _add_caption(s,
                 "6 nodes call an LLM. 3 are pure rules / local embeddings / SQLite. Every call logged.",
                 top=1.05)
    _add_bullets(s, [
        "•  LLM-driven nodes: extract, classify, decide, draft (DeepSeek V4-Flash + R1 for reasoning)",
        "•  Local-embedding node: retrieval over the policy library (bge-large-en-v1.5)",
        "•  Pure-rules nodes: HITL router (deterministic 16-action table), executor recipe map",
        "•  SQLite-backed: HITL queue, pipeline runs, audit log (Postgres swap is a db_url change)",
        "•  Every LLM call logged: timestamp, task, model, tokens, cost, latency, case_id",
        "•  Pipeline is plain async function; LangGraph wrap deferred until pilot needs durable state",
    ], top=1.7, size=16)
    _add_callout_box(s,
                     "Cost discipline isn't a slogan — it's instrumented. logs/llm_calls.jsonl is the ledger.",
                     top=5.9, height=0.9)
    _add_footer(s, 7)

    # ============== Slide 8 — Model + cost discipline ==============
    s = pres.slides.add_slide(blank)
    _add_brand_strip(s, "Model strategy + cost discipline")
    _add_caption(s, "Open-source first. Closed-model fallback per task. Local embeddings.",
                 top=1.05)
    _add_bullets(s, [
        "•  Default per task: DeepSeek V4-Flash (extraction, classification, drafting)",
        "•  Reasoning: DeepSeek R1 (decision-support — one call per invoice, ~$0.005)",
        "•  Embeddings: bge-large-en-v1.5 — local, free at runtime",
        "•  Closed-model fallback (Anthropic / OpenAI) wired but unused in build",
        "•  90–95% cost savings vs Claude Sonnet / GPT-4o at agent quality",
    ], top=1.7, size=17)

    # Stats row
    stats = [
        ("$3.74", "Total LLM spend\nentire build to date"),
        ("2,482", "LLM calls\nacross all tests + demos"),
        ("$0.0015", "Average cost\nper LLM call"),
        ("$0.007", "Cost per invoice\nend-to-end"),
    ]
    x = 0.7
    for value, label in stats:
        tile = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(5.0), Inches(3.0), Inches(1.7),
        )
        tile.fill.solid()
        tile.fill.fore_color.rgb = SOFT
        tile.line.color.rgb = NAVY
        tile.line.width = Pt(0.5)
        tile.text_frame.margin_top = Inches(0.15)
        tile.text_frame.margin_left = Inches(0.2)
        tile.text_frame.margin_right = Inches(0.2)
        tx = tile.text_frame
        tx.text = value
        tx.paragraphs[0].alignment = PP_ALIGN.CENTER
        for run in tx.paragraphs[0].runs:
            run.font.name = FONT
            run.font.size = Pt(28)
            run.font.bold = True
            run.font.color.rgb = NAVY
        p = tx.add_paragraph()
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.name = FONT
            run.font.size = Pt(11)
            run.font.color.rgb = MUTED
        x += 3.15
    _add_footer(s, 8)

    # ============== Slide 9 — Stage 9: the recurring revenue moat ==============
    s = pres.slides.add_slide(blank)
    _add_brand_strip(s, "Stage 9 — the recurring revenue moat")
    _add_caption(s,
                 "Most consultants ship and disappear. We ship + measure + tune quarterly. That's the lock-in.",
                 top=1.05)
    _add_bullets(s, [
        "•  Every LLM call → cost ledger (logs/llm_calls.jsonl). Auditable.",
        "•  Every approved case → HITL audit trail. Auditable.",
        "•  Every pipeline run → SQLite snapshot with full per-node trace. Re-openable.",
        "•  Quarterly Stage 9 report to buyer: auto-pass rate trend, classification mix, cost per task, latency p95, policy snippet usage",
        "•  Built-in dashboard at /stage9 (you saw it). Real numbers, not a slide.",
    ], top=1.7, size=16)
    _add_callout_box(s,
                     "Stage 9 is what justifies the quarterly retainer. It's also what makes us "
                     "hard to displace once we're embedded in a buyer's quarterly review motion.",
                     top=5.6, height=1.2)
    _add_footer(s, 9)

    # ============== Slide 10 — The buyer pitch ==============
    s = pres.slides.add_slide(blank)
    _add_brand_strip(s, "The buyer pitch — what we say")
    _add_caption(s,
                 "Internal alignment matters. Here's the one-liner + the three reasons buyers buy.",
                 top=1.05)

    # One-liner callout
    _add_callout_box(s,
                     "\"We're the firm that says NO to AI use cases that don't pencil out — "
                     "and proves it with the process data, before we sell you the build.\"",
                     top=1.7, height=1.3, bg=SOFT)

    _add_bullets(s, [
        "•  Specificity — every buyer asks \"does this work with our SAP?\" We say yes + demo on tape.",
        "•  Stage 9 — recurring measurement, quarterly review. The differentiator from \"ship and disappear.\"",
        "•  Kill gate — we'll walk away from a $1M engagement if the data says it won't deliver.",
    ], top=3.5, size=17)

    _add_callout_box(s,
                     "Honest credibility beats slick promises. The kill gate is what earns the trust "
                     "that lets us close the bigger engagements.",
                     top=5.8, height=1.0, bg=SOFT)
    _add_footer(s, 10)

    # ============== Slide 11 — Pipeline + roadmap ==============
    s = pres.slides.add_slide(blank)
    _add_brand_strip(s, "Where we are + what's next")
    _add_caption(s,
                 "Honest timeline. Pipeline conversations in 90 days; first billable engagement signed in 120–150.",
                 top=1.05)
    _add_bullets(s, [
        "TODAY  →  9/9 logic nodes shipped. Mock executor live. Demo-ready.",
        "DAY 14  →  SAP credentials in hand. Real connector + real action executor.",
        "DAY 30  →  Full SAP integration recorded; buyer-facing demo ready.",
        "DAY 60  →  Three buyer conversations active.",
        "DAY 90  →  One engagement in proposal stage.",
        "DAY 120–150  →  First billable engagement signed.",
    ], top=1.8, size=18)
    _add_callout_box(s,
                     "Blocked today: SAP credentials (trial activated this week; setup walking through).  "
                     "Deferred until pilot: LangGraph wrap, Postgres migration, auth.",
                     top=5.6, height=1.3, bg=SOFT)
    _add_footer(s, 11)

    # ============== Slide 12 — The ask ==============
    s = pres.slides.add_slide(blank)
    _add_brand_strip(s, "The ask")
    _add_caption(s,
                 "Specific items per person. Q&A buffer follows this slide.",
                 top=1.05)

    asks = [
        ("Sridhar", "Air cover for the 120–150 day timeline. 2–3 warm intros to buyers in your network with known AP-exception pain. Push to reactivate Botminds + Covasant partnerships."),
        ("CTO", "Architecture review in next 2 weeks. Decide: LangGraph wrap now or at pilot? SQLite okay or Postgres now? Where to add more tests?"),
        ("Procurement Lead", "Identify 5–8 specific exception scenarios from your engagement experience worth adding to our golden case set. Sharpens the agent's real-world behavior."),
        ("Process Lead", "Co-author the buyer-facing version of this demo. Internal one is engineer-friendly; buyer one needs reviewer/auditor framing."),
    ]
    y = 1.7
    for who, ask in asks:
        # Name box
        box = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.7), Inches(y), Inches(2.6), Inches(1.0),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = NAVY
        box.line.fill.background()
        _set_text(box.text_frame, who, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Ask text
        tx = s.shapes.add_textbox(Inches(3.5), Inches(y + 0.1), Inches(9.5), Inches(0.9))
        _set_text(tx.text_frame, ask, size=14, color=TEXT)
        y += 1.2

    _add_footer(s, 12)

    return pres


def main() -> None:
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    pres = build()
    pres.save(str(OUTPUT))
    print(f"Wrote {OUTPUT}")
    print(f"  Slides: {len(pres.slides)}")


if __name__ == "__main__":
    main()
